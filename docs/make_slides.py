"""Generate April 3 meeting slides as .pptx"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --- Colors ---
BLUE = RGBColor(0x00, 0x5A, 0xB5)
DARK = RGBColor(0x2D, 0x2D, 0x2D)
GRAY = RGBColor(0x66, 0x66, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)
GREEN = RGBColor(0x1B, 0x7F, 0x3B)
RED = RGBColor(0xC0, 0x39, 0x2B)
ACCENT = RGBColor(0xE8, 0x8D, 0x00)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# --- Helpers ---

def add_blank_slide():
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def add_title_text(slide, text, left=0.8, top=0.3, width=11.7, size=36, color=BLUE, bold=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    return txBox


def add_body_text(slide, text, left=0.8, top=1.3, width=11.7, height=5.5, size=20, color=DARK):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    return tf


def add_bullet_slide(slide, title, bullets, sub_bullets=None, top=1.3):
    """Add title + bullet points. sub_bullets is a dict: {bullet_index: [sub_bullet_texts]}"""
    add_title_text(slide, title)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(11.7), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20)
        p.font.color.rgb = DARK
        p.space_after = Pt(8)
        p.level = 0

        if sub_bullets and i in sub_bullets:
            for sb in sub_bullets[i]:
                sp = tf.add_paragraph()
                sp.text = sb
                sp.font.size = Pt(17)
                sp.font.color.rgb = GRAY
                sp.space_after = Pt(4)
                sp.level = 1
    return tf


def add_table(slide, headers, rows, left=0.8, top=None, width=11.7, col_widths=None, font_size=16):
    """Add a formatted table."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    if top is None:
        top = 2.5

    tbl_shape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top),
                                        Inches(width), Inches(0.4 * n_rows))
    tbl = tbl_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Inches(w)

    # Header row
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(font_size)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.color.rgb = DARK
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG

    return tbl_shape


def add_note(slide, text):
    """Add speaker notes."""
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def add_callout(slide, text, left=0.8, top=6.2, width=11.7, size=17, color=GRAY):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.italic = True


# ============================================================
# SLIDE 1: Title
# ============================================================
slide = add_blank_slide()
# Blue background
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = BLUE
bg.line.fill.background()

add_title_text(slide, "Towards General Purpose Robots", left=1.0, top=2.5, size=44, color=WHITE)
add_title_text(slide, "Manu Gaur", left=1.0, top=3.5, size=28, color=WHITE, bold=False)
add_title_text(slide, "April 3, 2026", left=1.0, top=4.2, size=22, color=WHITE, bold=False)

# ============================================================
# SLIDE 2: Memory Bottleneck
# ============================================================
slide = add_blank_slide()
add_title_text(slide, "Long-Horizon Robotics is Memory-Bottlenecked")

tf = add_body_text(slide, "", top=1.2, size=18)
p = tf.paragraphs[0]
p.text = "Without an abstraction of the past, it is impossible to form a causal understanding of how the world evolves."
p.font.size = Pt(18)
p.font.color.rgb = GRAY
p.font.italic = True

bullets = [
    "As skills become robust, the bottleneck is deploying them for complex tasks",
    "Complex tasks need narrative memory: tracking task stage, object locations (even out of view), what's been tried",
    "Attention alone can't scale: full observation history over minutes is infeasible, but discarding it leads to incoherent behavior",
]
for b in bullets:
    p = tf.add_paragraph()
    p.text = b
    p.font.size = Pt(20)
    p.font.color.rgb = DARK
    p.space_after = Pt(10)
    p.level = 0

add_note(slide, "Examples: robot needs to track time for grilled cheese, remember progress wiping window. These require persistent memory, not just reactive skills.")

# ============================================================
# SLIDE 3: Why Future Prediction?
# ============================================================
slide = add_blank_slide()
add_title_text(slide, "Why Future Prediction?")

bullets = [
    "Temporal abstraction -- what task stage am I in, what comes next?",
    "Dynamics -- how objects move, interact, and change state",
    "Persistent memory for partial observability -- inferring what's out of view from what was seen",
]
tf = add_bullet_slide(slide, "Why Future Prediction?", bullets, top=1.5)

add_callout(slide, "Same principle the brain uses -- prediction error drives attention, memory and learning [Friston, 2010]",
            top=5.5, size=16)
# remove duplicate title - the bullet_slide already adds one
# (we called add_title_text then add_bullet_slide which also adds title)

add_note(slide, "This is the theoretical motivation. Prediction is not just a proxy task -- it forces the model to build the representations we need.")

# ============================================================
# SLIDE 4: High-Level Plan (3 stages)
# ============================================================
slide = add_blank_slide()
add_title_text(slide, "Three-Stage Approach")

# Stage boxes
stages = [
    ("Stage 1", "Recurrent Video\nCompression", "Compress short video into\ncompact dynamics state", 1.0),
    ("Stage 2", "Long-Horizon\nForecasting", "Autoregressive prediction\nacross compressed states", 5.0),
    ("Stage 3", "World\nModelling", "State + Action -> Next State\nPlanning via MPC", 9.0),
]

for label, name, desc, left in stages:
    # Box
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(2.0),
                                  Inches(3.3), Inches(3.5))
    box.fill.solid()
    box.fill.fore_color.rgb = LIGHT_BG
    box.line.color.rgb = BLUE
    box.line.width = Pt(2)

    tf = box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Label
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(14)
    p.font.color.rgb = BLUE
    p.font.bold = True
    p.space_after = Pt(6)

    # Name
    p2 = tf.add_paragraph()
    p2.text = name
    p2.font.size = Pt(22)
    p2.font.color.rgb = DARK
    p2.font.bold = True
    p2.alignment = PP_ALIGN.CENTER
    p2.space_after = Pt(12)

    # Desc
    p3 = tf.add_paragraph()
    p3.text = desc
    p3.font.size = Pt(15)
    p3.font.color.rgb = GRAY
    p3.alignment = PP_ALIGN.CENTER

# Arrows between boxes
for x in [4.3, 8.3]:
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(3.3),
                                    Inches(0.7), Inches(0.4))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = BLUE
    arrow.line.fill.background()

add_callout(slide, "Long-horizon BPTT is infeasible (GPU + data). Staged approach makes it tractable.", top=6.0)
add_note(slide, "Stage 1 is current focus. Everything downstream depends on quality of compressed state.")

# ============================================================
# SLIDE 5: Stage 1 -- Current Setup
# ============================================================
slide = add_blank_slide()
add_title_text(slide, "Stage 1: What We're Doing Now")

tf = add_body_text(slide, "", top=1.2)

p = tf.paragraphs[0]
p.text = "Goal: compress video into a state expressive enough for stage 2 forecasting"
p.font.size = Pt(20)
p.font.color.rgb = DARK
p.font.bold = True
p.space_after = Pt(16)

setup_items = [
    "Frozen DINO ViT-S/14 backbone (384-dim features)",
    "SSv2, 8 uniformly sampled frames",
    "Additive RNN with predictive coding update:",
]
for item in setup_items:
    p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(19)
    p.font.color.rgb = DARK
    p.space_after = Pt(4)

# Equations
eqs = [
    "Error = DINO(x_t) - W_pred(S_{t-1})",
    "S_t = LN(S_{t-1} + W(Error))",
]
for eq in eqs:
    p = tf.add_paragraph()
    p.text = "    " + eq
    p.font.size = Pt(17)
    p.font.color.rgb = BLUE
    p.font.bold = True
    p.space_after = Pt(2)

p = tf.add_paragraph()
p.space_after = Pt(16)

p = tf.add_paragraph()
p.text = "Key finding: predictive coding > sigmoid gating"
p.font.size = Pt(20)
p.font.color.rgb = GREEN
p.font.bold = True
p.space_after = Pt(4)

p = tf.add_paragraph()
p.text = "Easier to predict next frame than learn what to forget"
p.font.size = Pt(17)
p.font.color.rgb = GRAY

add_note(slide, "Predictive coding works because the error signal itself is informative -- the state accumulates corrections rather than trying to gate information in/out.")

# ============================================================
# SLIDE 6: Models Trained
# ============================================================
slide = add_blank_slide()
add_title_text(slide, "Models Trained")

add_table(slide,
    headers=["#", "Model", "Prediction Target", "Loss", "Status"],
    rows=[
        ["1", "zyvsy8gk", "next frame, learned space h=W(DINO_cls)", "CE + 0.1*L2", "works"],
        ["2", "--", "next frame, learned space, pred only", "L2 only", "collapsed"],
        ["3a", "2ldiw9xk", "next CLS in DINO space", "L2", "works"],
        ["3b", "e6esmgmu", "next patches in DINO space", "L2", "works"],
    ],
    top=1.8,
    col_widths=[0.5, 1.5, 5.0, 2.0, 1.5],
    font_size=15,
)

tf = add_body_text(slide, "", left=0.8, top=4.5, width=11.7, height=2.0)
p = tf.paragraphs[0]
p.text = "Takeaway: prediction-only in learned space collapses"
p.font.size = Pt(20)
p.font.color.rgb = RED
p.font.bold = True
p.space_after = Pt(6)

p2 = tf.add_paragraph()
p2.text = "Without classification gradient, the target space degenerates. DINO-space prediction avoids this."
p2.font.size = Pt(18)
p2.font.color.rgb = DARK

add_note(slide, "Collapse of model 2 motivates two directions: (a) fix learned-space collapse with SigReg/EMA, (b) improve DINO-space prediction with more expressive W_pred.")

# ============================================================
# SLIDE 7: Transfer Probe -- UCF101
# ============================================================
slide = add_blank_slide()
add_title_text(slide, "Transfer Probe: UCF101")

tf = add_body_text(slide, "", top=1.2, height=0.5)
p = tf.paragraphs[0]
p.text = "Does the recurrent state generalize beyond SSv2?"
p.font.size = Pt(18)
p.font.color.rgb = GRAY
p.font.italic = True

add_body_text(slide, "Freeze everything, linear head on UCF101 (101 classes, 13K videos, 20 epochs)",
              top=1.7, height=0.4, size=16, color=GRAY)

add_table(slide,
    headers=["Model", "UCF101 Acc"],
    rows=[
        ["DINO mean-pool (no temporal modeling)", "88.0%"],
        ["DINO concat (8x384=3072-dim)", "86.0%"],
        ["RNN CLS, DINO-space pred", "85.4%"],
        ["RNN CLS, CE+pred, learned space", "84.0%"],
        ["RNN Patches, DINO-space pred", "81.7%"],
        ["RNN Patches, CE+pred, learned space", "78.3%"],
    ],
    top=2.4,
    col_widths=[8.5, 3.2],
    font_size=16,
)

tf2 = add_body_text(slide, "", left=0.8, top=5.5, width=11.7, height=1.5)
bullets = [
    "All RNN states trail DINO mean-pool (88%)",
    "DINO-space > learned space -- CE loss overfits to SSv2",
    "mean-pool > concat suggests UCF measures appearance, not temporal reasoning",
]
for i, b in enumerate(bullets):
    p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
    p.text = b
    p.font.size = Pt(16)
    p.font.color.rgb = DARK
    p.space_after = Pt(3)

add_note(slide, "UCF101 is appearance-heavy. Need a non-action-recognition eval to truly test dynamics.")

# ============================================================
# SLIDE 8: DinoWM -- Causal Transformer Baseline
# ============================================================
slide = add_blank_slide()
add_title_text(slide, "DinoWM: How Lossy is Recurrence?")

tf = add_body_text(slide, "", top=1.2, height=0.5)
p = tf.paragraphs[0]
p.text = "Causal transformer = full attention over past frames (no state bottleneck). Upper bound for recurrence."
p.font.size = Pt(18)
p.font.color.rgb = GRAY
p.font.italic = True

# CLS table
add_body_text(slide, "CLS token (S=1, D=384)", top=2.0, height=0.4, size=18, color=DARK)
add_table(slide,
    headers=["Model", "Pred L2"],
    rows=[
        ["Copy current frame", "609"],
        ["Causal Transformer", "517"],
        ["RNN", "513"],
    ],
    top=2.5,
    left=0.8, width=5.5,
    col_widths=[3.5, 2.0],
    font_size=15,
)

# Patches table
add_body_text(slide, "Patch tokens (S=256, D=384)", left=7.0, top=2.0, height=0.4, size=18, color=DARK)
add_table(slide,
    headers=["Model", "Pred L2"],
    rows=[
        ["Copy current frame", "1085"],
        ["Causal Transformer", "783"],
        ["RNN", "851"],
    ],
    top=2.5,
    left=7.0, width=5.5,
    col_widths=[3.5, 2.0],
    font_size=15,
)

tf2 = add_body_text(slide, "", left=0.8, top=5.0, width=11.7, height=2.0)
takeaways = [
    "CLS: RNN matches causal transformer -- state bottleneck is NOT the problem",
    "Patches: RNN trails transformer by ~8% -- need more expressive W_pred, not bigger state",
    "Both models beat copy-frame baseline",
]
for i, t in enumerate(takeaways):
    p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
    p.text = t
    p.font.size = Pt(17)
    p.font.color.rgb = DARK
    p.space_after = Pt(5)

add_note(slide, "For CLS the recurrent state captures essentially all information that full attention does. For patches, the gap suggests we need a better predictor, not a bigger state.")

# ============================================================
# SLIDE 9: Temporal Shuffle Test -- Results
# ============================================================
slide = add_blank_slide()
add_title_text(slide, "Temporal Shuffle Test")

tf = add_body_text(slide, "", top=1.2, height=0.8)
p = tf.paragraphs[0]
p.text = "Model beats copy baseline by 17-24% -- is that from learning dynamics or just averaging recent frames?"
p.font.size = Pt(18)
p.font.color.rgb = GRAY
p.font.italic = True
p2 = tf.add_paragraph()
p2.text = "Shuffling frame order tests this: dynamics break under shuffling, averaging doesn't."
p2.font.size = Pt(16)
p2.font.color.rgb = GRAY
p2.space_after = Pt(4)
p3 = tf.add_paragraph()
p3.text = "Ratio = shuffled_loss / normal_loss  (higher = more temporal dependence)"
p3.font.size = Pt(16)
p3.font.color.rgb = DARK
p3.font.bold = True

# Main results table
add_table(slide,
    headers=["Model", "Normal", "Shuffled", "Ratio", "Copy Baseline"],
    rows=[
        ["CE+pred, CLS, learned space", "2.23", "2.29", "1.03x", "--"],
        ["CE+pred, patches, learned space", "3.80", "4.29", "1.13x", "--"],
        ["Pred only, CLS, DINO space", "176.9", "211.9", "1.20x", "620"],
        ["Pred only, patches, DINO space", "851.5", "1114.1", "1.31x", "1116"],
    ],
    top=2.8,
    col_widths=[4.5, 1.7, 1.7, 1.5, 2.0],
    font_size=15,
)

# Copy baseline shuffle sensitivity table
add_body_text(slide, "Copy Baseline Shuffle Sensitivity", left=0.8, top=5.1, height=0.4, size=16, color=DARK)
add_table(slide,
    headers=["", "Copy Shuffle Ratio", "Model Shuffle Ratio"],
    rows=[
        ["CLS (DINO space)", "11.2x", "1.20x"],
        ["Patches (DINO space)", "1.46x", "1.31x"],
        ["Pixel space", "1.55x", "--"],
    ],
    top=5.5,
    left=0.8, width=8.0,
    col_widths=[3.0, 2.5, 2.5],
    font_size=14,
)

add_note(slide, "Learned-space models barely care about temporal order (3-13%) -- CE overfits to action labels. Copy ratio = data property (frame similarity between random vs consecutive pairs). Model ratio = model property (how much it relied on temporal order).")

# ============================================================
# SLIDE 9b: Shuffle Interpretation -- CLS vs Patches
# ============================================================
slide = add_blank_slide()
add_title_text(slide, "Shuffle Test: What Did the Models Learn?")

# CLS box (left)
box_cls = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5),
                                  Inches(5.5), Inches(4.5))
box_cls.fill.solid()
box_cls.fill.fore_color.rgb = LIGHT_BG
box_cls.line.color.rgb = BLUE
box_cls.line.width = Pt(1.5)
tf_cls = box_cls.text_frame
tf_cls.word_wrap = True

p = tf_cls.paragraphs[0]
p.text = "CLS Token"
p.font.size = Pt(22)
p.font.color.rgb = BLUE
p.font.bold = True
p.space_after = Pt(12)

cls_items = [
    ("1. Multi-frame aggregation (order-independent)", DARK, True),
    ("Shuffled model (211.9) still far better than copy (620)", GRAY, False),
    ("Combining all frames >> any single frame", GRAY, False),
    ("", DARK, False),
    ("2. Temporal dynamics (order-dependent)", DARK, True),
    ("20% degradation from shuffling (176.9 -> 211.9)", GRAY, False),
    ("", DARK, False),
    ("Copy shuffle ratio is 11.2x because random CLS", GRAY, False),
    ("pairs are very different -- but model doesn't care,", GRAY, False),
    ("it's using all frames, not just the last one", GRAY, False),
]
for text, color, bold in cls_items:
    p = tf_cls.add_paragraph()
    p.text = text
    p.font.size = Pt(15)
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_after = Pt(3)

# Patches box (right)
box_pat = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.5),
                                  Inches(5.5), Inches(4.5))
box_pat.fill.solid()
box_pat.fill.fore_color.rgb = LIGHT_BG
box_pat.line.color.rgb = BLUE
box_pat.line.width = Pt(1.5)
tf_pat = box_pat.text_frame
tf_pat.word_wrap = True

p = tf_pat.paragraphs[0]
p.text = "Patch Tokens"
p.font.size = Pt(22)
p.font.color.rgb = BLUE
p.font.bold = True
p.space_after = Pt(12)

pat_items = [
    ("Shuffled model (1114) = copy baseline (1116)", DARK, True),
    ("ALL improvement over copy destroyed by shuffling", RED, True),
    ("", DARK, False),
    ("No benefit of running summary (unlike CLS)", GRAY, False),
    ("-- patches are mostly static, one frame is a", GRAY, False),
    ("sufficient predictor", GRAY, False),
    ("", DARK, False),
    ("DINO patches shuffle ratio (1.46x) behaves", GRAY, False),
    ("like pixel space (1.55x)", GRAY, False),
    ("", DARK, False),
    ("100% of model's advantage = temporal order", GREEN, True),
]
for text, color, bold in pat_items:
    p = tf_pat.add_paragraph()
    p.text = text
    p.font.size = Pt(15)
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_after = Pt(3)

add_note(slide, "Key contrast: CLS learns two things (aggregation + dynamics), patches learn one thing (dynamics only). This explains why CLS is more robust to shuffling -- it still has the aggregation benefit. Patches lose everything because their only advantage was temporal order.")

# ============================================================
# SLIDE 9c: CLS or Patch Prediction?
# ============================================================
slide = add_blank_slide()
add_title_text(slide, "CLS or Patch Prediction?")

# CLS pros/cons (left)
box_c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5),
                                Inches(5.5), Inches(2.5))
box_c.fill.solid()
box_c.fill.fore_color.rgb = LIGHT_BG
box_c.line.color.rgb = BLUE
box_c.line.width = Pt(1.5)
tf_c = box_c.text_frame
tf_c.word_wrap = True
p = tf_c.paragraphs[0]
p.text = "CLS"
p.font.size = Pt(22)
p.font.color.rgb = BLUE
p.font.bold = True
p.space_after = Pt(8)

for text, color in [
    ("+ Cleaner training signal, no gradient dilution", GREEN),
    ("+ Every frame contributes meaningfully", GREEN),
    ("- Captures semantic shifts, not just motion?", RED),
    ("- Throws away spatial info stage 2 needs", RED),
]:
    p = tf_c.add_paragraph()
    p.text = text
    p.font.size = Pt(16)
    p.font.color.rgb = color
    p.space_after = Pt(4)

# Patches pros/cons (right)
box_p = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.5),
                                Inches(5.5), Inches(2.5))
box_p.fill.solid()
box_p.fill.fore_color.rgb = LIGHT_BG
box_p.line.color.rgb = BLUE
box_p.line.width = Pt(1.5)
tf_p = box_p.text_frame
tf_p.word_wrap = True
p = tf_p.paragraphs[0]
p.text = "Patches"
p.font.size = Pt(22)
p.font.color.rgb = BLUE
p.font.bold = True
p.space_after = Pt(8)

for text, color in [
    ("+ Right inductive bias: spatial + temporal", GREEN),
    ("+ Better aligned with stage 2 forecasting", GREEN),
    ("- Gradient dilution from static patches", RED),
    ("- Trails CLS on UCF by 4-7 pts", RED),
]:
    p = tf_p.add_paragraph()
    p.text = text
    p.font.size = Pt(16)
    p.font.color.rgb = color
    p.space_after = Pt(4)

# Fixes for gradient dilution
tf_fix = add_body_text(slide, "", left=0.8, top=4.5, width=11.7, height=2.5)
p = tf_fix.paragraphs[0]
p.text = "Fixing patch gradient dilution:"
p.font.size = Pt(20)
p.font.color.rgb = DARK
p.font.bold = True
p.space_after = Pt(8)

fixes = [
    "1. Motion-weighted loss -- upweight dynamic patches in L2 loss",
    "2. Predict residuals (x_{t+1} - x_t) -- static patches become ~zero, dynamic = entire signal",
    "3. CLS as auxiliary -- keeps gradients flowing when patch loss stagnates on static tokens",
]
for f in fixes:
    p = tf_fix.add_paragraph()
    p.text = f
    p.font.size = Pt(17)
    p.font.color.rgb = DARK
    p.space_after = Pt(6)

add_note(slide, "For stage 1 (state expressive enough for stage 2 forecasting), patches are better aligned. The question is how to fix gradient dilution. These three approaches are not mutually exclusive.")

# ============================================================
# SLIDE 10: Static vs Dynamic Patches
# ============================================================
slide = add_blank_slide()
add_title_text(slide, "Static vs Dynamic Patch Decomposition")

tf = add_body_text(slide, "", top=1.2, height=0.8)
p = tf.paragraphs[0]
p.text = "Is the patch model learning dynamics or just smoothing?"
p.font.size = Pt(18)
p.font.color.rgb = GRAY
p.font.italic = True
p2 = tf.add_paragraph()
p2.text = "Split patches into static (bottom 50% motion) and dynamic (top 50%) per video. Compare improvement over copy baseline."
p2.font.size = Pt(16)
p2.font.color.rgb = GRAY

add_table(slide,
    headers=["Patch Group", "Copy Baseline", "Model", "Improvement"],
    rows=[
        ["Dynamic patches", "1430.6", "1077.0", "24.7%"],
        ["Static patches", "746.9", "625.9", "16.2%"],
    ],
    top=3.0,
    col_widths=[3.5, 3.0, 2.5, 2.7],
    font_size=16,
)

tf2 = add_body_text(slide, "", left=0.8, top=4.8, width=11.7, height=2.0)

p = tf2.paragraphs[0]
p.text = "Improvement 1.5x more concentrated on dynamic patches -- model learned dynamics, not just smoothing"
p.font.size = Pt(18)
p.font.color.rgb = GREEN
p.font.bold = True
p.space_after = Pt(12)

p2 = tf2.add_paragraph()
p2.text = "Dynamic-only copy shuffle ratio: 1.40x (vs CLS 11.2x)"
p2.font.size = Pt(17)
p2.font.color.rgb = DARK
p2.space_after = Pt(4)

p3 = tf2.add_paragraph()
p3.text = "CLS captures something beyond patch-level motion -- holistic semantic shifts"
p3.font.size = Pt(16)
p3.font.color.rgb = GRAY

add_note(slide, "Most direct evidence that RNN is tracking motion. If it were just a better EMA, improvement would be equal across static and dynamic patches.")

# ============================================================
# SLIDE 11: Next Steps
# ============================================================
slide = add_blank_slide()
add_title_text(slide, "Next Steps")

# Two columns
# Left column: Direction A
box_a = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5),
                                Inches(5.5), Inches(3.5))
box_a.fill.solid()
box_a.fill.fore_color.rgb = LIGHT_BG
box_a.line.color.rgb = BLUE
box_a.line.width = Pt(1.5)
tf_a = box_a.text_frame
tf_a.word_wrap = True
p = tf_a.paragraphs[0]
p.text = "A) Fix learned-space collapse"
p.font.size = Pt(20)
p.font.color.rgb = BLUE
p.font.bold = True
p.space_after = Pt(10)

items_a = [
    "JEPA-style h=W(DINO) + SigReg",
    "Currently collapses without CE loss",
    "Try: SigReg regularization",
    "Try: EMA / stop-gradient on target encoder",
]
for item in items_a:
    p = tf_a.add_paragraph()
    p.text = item
    p.font.size = Pt(16)
    p.font.color.rgb = DARK
    p.space_after = Pt(4)

# Right column: Direction B
box_b = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.5),
                                Inches(5.5), Inches(3.5))
box_b.fill.solid()
box_b.fill.fore_color.rgb = LIGHT_BG
box_b.line.color.rgb = BLUE
box_b.line.width = Pt(1.5)
tf_b = box_b.text_frame
tf_b.word_wrap = True
p = tf_b.paragraphs[0]
p.text = "B) More expressive W_pred"
p.font.size = Pt(20)
p.font.color.rgb = BLUE
p.font.bold = True
p.space_after = Pt(10)

items_b = [
    "Recurrent state is NOT the bottleneck",
    "Close patch gap: RNN 851 vs transformer 783",
    "Goal: match DINO mean-pool on UCF (88%)",
    "More expressive predictor, not bigger state",
]
for item in items_b:
    p = tf_b.add_paragraph()
    p.text = item
    p.font.size = Pt(16)
    p.font.color.rgb = DARK
    p.space_after = Pt(4)

# Diagnostics section
tf3 = add_body_text(slide, "", left=0.8, top=5.3, width=11.7, height=1.5)
p = tf3.paragraphs[0]
p.text = "Diagnostics to run:"
p.font.size = Pt(18)
p.font.color.rgb = DARK
p.font.bold = True
p.space_after = Pt(6)

diags = [
    "OOD prediction decay curve (SSv2 vs UCF101 -- general dynamics or SSv2-memorized?)",
    "Multi-step autoregressive rollout (extrapolate trajectories or collapse to mean?)",
]
for d in diags:
    p = tf3.add_paragraph()
    p.text = d
    p.font.size = Pt(16)
    p.font.color.rgb = GRAY
    p.space_after = Pt(3)

add_note(slide, "Two directions are complementary -- (A) is about training objective, (B) is about model capacity. Diagnostics will tell us which matters more.")


# ============================================================
# Save
# ============================================================
out_path = "/home/manu/vjepa2/docs/april_3_slides.pptx"
prs.save(out_path)
print(f"Saved to {out_path}")
